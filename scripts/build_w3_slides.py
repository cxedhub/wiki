#!/usr/bin/env python3
"""Build Workshop-3-Slides.pptx to match slides-w3.html (the canonical web deck).

Regenerated whenever the live agenda or web slides change so the downloadable
deck at static/pds/craft-pd-series/Workshop-3-Slides.pptx stays in sync.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Palette mirrors slides-w3.html
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)
LIGHT_BG = RGBColor(0xFA, 0xFB, 0xFE)
TEXT_DARK = RGBColor(0x1B, 0x2A, 0x4A)
TEXT_MUTED = RGBColor(0x5A, 0x7F, 0xA0)
TEXT_BODY = RGBColor(0x2C, 0x4A, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

C_CONTEXT = RGBColor(0xD4, 0xA8, 0x43)   # yellow
C_REFRAME = RGBColor(0xC0, 0x39, 0x2B)   # red
C_ASSEMBLE = RGBColor(0x27, 0xAE, 0x60)  # green
C_FORTIFY = RGBColor(0x1A, 0xBC, 0x9C)   # teal
C_TRANSFER = RGBColor(0x8E, 0x6B, 0xBF)  # purple
C_BREAK = RGBColor(0x99, 0x99, 0x99)
C_LINK = RGBColor(0x1A, 0xBC, 0x9C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, font="DM Sans", size=14, bold=False,
             italic=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT_DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = "DM Sans"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_badge(slide, x, y, text, fill, *, text_color=WHITE, outline=False,
              width=Inches(1.1), height=Inches(0.32)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
    if outline:
        shp.fill.background()
        shp.line.color.rgb = fill
        shp.line.width = Pt(1.25)
        color = fill
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.fill.background()
        color = text_color
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "DM Sans"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = color
    return shp


def dark_bg(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)


def light_bg(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_BG)


def top_bar(slide, color):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), color)


def bottom_bar(slide, color):
    add_rect(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), color)


def left_accent(slide, color):
    add_rect(slide, 0, 0, Inches(0.08), SLIDE_H, color)


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)
    top_bar(s, C_ASSEMBLE)
    add_text(s, Inches(0.9), Inches(0.6), Inches(9), Inches(0.4),
             "WORKSHOP 03  ·  APRIL 25, 2026",
             font="JetBrains Mono", size=11, color=TEXT_MUTED)
    add_text(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(2.2),
             "Programming Edge/IoT\nSystems with AI",
             font="DM Serif Display", size=48, color=WHITE)
    add_text(s, Inches(0.9), Inches(4.2), Inches(10), Inches(1.2),
             "A two-node IoT system built entirely in the MakeCode simulator — "
             "sensor node talks to aggregator node, LLM as co-pilot. "
             "Your kit ships after.",
             size=16, color=TEXT_MUTED)
    add_text(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.8),
             "April 25, 2026  ·  8:30 AM – 12:00 PM PST  ·  Virtual\n"
             "CRAFT PD Series  ·  UCF DRACO Lab & School of Teacher Education",
             size=10, color=TEXT_MUTED)


def slide_craft_cycle(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(s)
    add_text(s, Inches(0.9), Inches(0.5), Inches(11), Inches(0.9),
             "The CRAFT Cycle",
             font="DM Serif Display", size=36, color=TEXT_DARK)
    add_text(s, Inches(0.9), Inches(1.4), Inches(11), Inches(0.6),
             "Today's session IS a CRAFT cycle. "
             "You'll experience it — then learn to build your own.",
             size=14, color=TEXT_MUTED)

    phases = [
        ("C", "Contextualize", "Why it matters — real-world framing", C_CONTEXT),
        ("R", "Reframe", "Surface the wrong model → install the right one", C_REFRAME),
        ("A", "Assemble", "I Do → We Do → You Do", C_ASSEMBLE),
        ("F", "Fortify", "Verify with tools + AI; errors are features", C_FORTIFY),
        ("T", "Transfer", "Connect forward to your classroom", C_TRANSFER),
    ]
    gap = Inches(0.15)
    left = Inches(0.9)
    total_w = Inches(11.5)
    card_w = Emu((total_w - gap * (len(phases) - 1)) // len(phases))
    y = Inches(2.5)
    card_h = Inches(4.0)
    for i, (letter, name, desc, fill) in enumerate(phases):
        x = left + (card_w + gap) * i
        add_rect(s, x, y, card_w, card_h, fill)
        add_text(s, x, y + Inches(0.4), card_w, Inches(1.6), letter,
                 font="DM Serif Display", size=54, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(2.0), card_w, Inches(0.4), name.upper(),
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(2.6), card_w - Inches(0.3),
                 Inches(1.2), desc,
                 size=10, color=WHITE, align=PP_ALIGN.CENTER)


def content_slide(prs, *, accent, do_label, phase_label, phase_color,
                  time_label, title, body, bullets=None, link=None):
    """A 'light' content slide matching the HTML template."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(s)
    left_accent(s, accent)

    # Badges row
    badge_y = Inches(0.55)
    x = Inches(0.9)
    add_badge(s, x, badge_y, do_label, C_ASSEMBLE)
    x += Inches(1.2)
    if phase_label:
        add_badge(s, x, badge_y, phase_label, phase_color, outline=True,
                  width=Inches(1.5))
        x += Inches(1.6)
    add_text(s, Inches(10.5), badge_y, Inches(2.2), Inches(0.32),
             time_label, font="JetBrains Mono", size=11,
             color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.9), Inches(1.1), Inches(11.5), Inches(1.1),
             title, font="DM Serif Display", size=30, color=TEXT_DARK)

    add_text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.0),
             body, size=15, color=TEXT_BODY)

    if bullets:
        add_bullets(s, Inches(1.1), Inches(3.1), Inches(11), Inches(4.0),
                    bullets, size=13, color=TEXT_DARK)

    if link:
        add_text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5),
                 link, size=13, bold=True, color=C_LINK)


def phase_intro_slide(prs, *, letter, title, subtitle, color, time_label):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)
    top_bar(s, color)
    add_text(s, Inches(11.0), Inches(0.5), Inches(1.8), Inches(0.4),
             time_label, font="JetBrains Mono", size=11,
             color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.RIGHT)
    add_text(s, Inches(0.9), Inches(3.2), Inches(6), Inches(2.2), letter,
             font="DM Serif Display", size=110, color=color)
    add_text(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.9), title,
             font="DM Serif Display", size=34, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.9), subtitle,
             size=14, color=TEXT_MUTED)


def break_slide(prs, *, time_label, title, body):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(s)
    left_accent(s, C_BREAK)
    add_badge(s, Inches(0.9), Inches(0.55), "BREAK", C_BREAK)
    add_text(s, Inches(10.5), Inches(0.55), Inches(2.2), Inches(0.32),
             time_label, font="JetBrains Mono", size=11,
             color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(1.1), Inches(11.5), Inches(1.0),
             title, font="DM Serif Display", size=30, color=TEXT_DARK)
    add_text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.5),
             body, size=15, color=TEXT_BODY)


def level_up_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(s)
    left_accent(s, C_ASSEMBLE)
    add_badge(s, Inches(0.9), Inches(0.55), "YOU DO", C_ASSEMBLE)
    add_text(s, Inches(10.5), Inches(0.55), Inches(2.2), Inches(0.32),
             "11:00 (30 min)", font="JetBrains Mono", size=11,
             color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(1.1), Inches(11.5), Inches(1.0),
             "Level Up: Choose Your Track",
             font="DM Serif Display", size=30, color=TEXT_DARK)

    col_y = Inches(2.4)
    col_h = Inches(4.2)
    col_w = Inches(5.6)
    add_rect(s, Inches(0.9), col_y, col_w, col_h,
             RGBColor(0xE8, 0xF8, 0xF0))
    add_text(s, Inches(1.2), col_y + Inches(0.3), col_w - Inches(0.6),
             Inches(0.5), "TRACK A: MULTI-SENSOR MESH",
             size=11, bold=True, color=TEXT_DARK)
    add_text(s, Inches(1.2), col_y + Inches(0.9), col_w - Inches(0.6),
             col_h - Inches(1.0),
             "Open your Sensor project in a 2nd browser tab for a 2nd Sensor "
             "Node (MakeCode pairs two sims per tab; tabs share the radio "
             "channel). Each sender includes an ID. Aggregator tracks "
             "per-sender averages and flags drop-offs.",
             size=12, color=TEXT_BODY)

    add_rect(s, Inches(6.8), col_y, col_w, col_h,
             RGBColor(0xF3, 0xED, 0xF9))
    add_text(s, Inches(7.1), col_y + Inches(0.3), col_w - Inches(0.6),
             Inches(0.5), "TRACK B: MICROPYTHON PORT",
             size=11, bold=True, color=TEXT_DARK)
    add_text(s, Inches(7.1), col_y + Inches(0.9), col_w - Inches(0.6),
             col_h - Inches(1.0),
             "Port sender + receiver to python.microbit.org. Ask an LLM for "
             "the translation; run both in the MicroPython simulator and "
             "verify messages still flow.",
             size=12, color=TEXT_BODY)


def resource_kit_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(s)
    top_bar(s, C_ASSEMBLE)
    add_text(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.9),
             "Your Resource Kit",
             font="DM Serif Display", size=32, color=TEXT_DARK)
    add_text(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.6),
             "Everything is yours. Bookmark, download, remix. "
             "Hardware ships after.",
             size=14, color=TEXT_MUTED)

    items = [
        ("Sensor + Aggregator .hex",
         "MakeCode files from today — flash when kit arrives"),
        ("MicroPython Radio Starters",
         "Sender + receiver in Python"),
        ("Sensor Reference",
         "All micro:bit V2 sensors"),
        ("NGSS Crosswalk",
         "IoT ↔ standards mapping"),
        ("Two-Node Lesson Template",
         "NGSS-aligned, customizable"),
        ("micro:bit V2 Kit",
         "Mailed after post-survey + lesson draft"),
    ]
    card_w = Inches(5.6)
    card_h = Inches(1.3)
    gap_y = Inches(0.2)
    gap_x = Inches(0.2)
    start_x = Inches(0.9)
    start_y = Inches(2.2)
    for i, (name, desc) in enumerate(items):
        row, col = divmod(i, 2)
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_rect(s, x, y, card_w, card_h, WHITE,
                 line=RGBColor(0xE8, 0xE8, 0xE8))
        add_text(s, x + Inches(0.25), y + Inches(0.2), card_w - Inches(0.4),
                 Inches(0.4), name, size=13, bold=True, color=TEXT_DARK)
        add_text(s, x + Inches(0.25), y + Inches(0.65), card_w - Inches(0.4),
                 Inches(0.6), desc, size=11, color=RGBColor(0x66, 0x66, 0x66))


def closing_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)
    bottom_bar(s, C_ASSEMBLE)
    add_text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(3.5),
             "You just built a two-node IoT system\n"
             "in a browser, with AI as your co-pilot,\n"
             "and verified it like an engineer.\n"
             "Your students can do this too —\n"
             "even in schools with zero hardware budget.",
             font="DM Serif Display", size=26, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.7),
             "CRAFT PD Series  ·  UCF DRACO Lab & School of Teacher Education\n"
             "Dr. Mike Borowczak & Dr. Andrea Borowczak",
             size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover
    slide_cover(prs)

    # 2. CRAFT cycle
    slide_craft_cycle(prs)

    # 3. 8:30 Simulator Check
    content_slide(
        prs,
        accent=RGBColor(0x5A, 0x7F, 0xA0),
        do_label="YOU DO", phase_label=None, phase_color=None,
        time_label="8:30 (15 min)",
        title="Simulator Check: Your micro:bit Is a Browser Tab",
        body="No hardware needed today.",
        bullets=[
            "Open makecode.microbit.org",
            "Drop show string \"hi\" in on start — watch the sim scroll",
            "Drop radio send number in — a second simulator appears",
            "That's our IoT system. Complete the pre-survey.",
        ],
        link="→ makecode.microbit.org",
    )

    # 4. Contextualize intro
    phase_intro_slide(
        prs, letter="C", title="Contextualize",
        subtitle="IoT is a conversation between two roles: "
                 "sensor nodes + aggregator nodes",
        color=C_CONTEXT, time_label="8:45",
    )

    # 5. 8:45 Brainstorm
    content_slide(
        prs, accent=C_CONTEXT,
        do_label="YOU DO", phase_label="CONTEXTUALIZE",
        phase_color=C_CONTEXT, time_label="8:45 (15 min)",
        title="Brainstorm: Two-Node Scenarios in Your Subject",
        body="In the shared doc:",
        bullets=[
            "Name a real-world sensor → aggregator pair from a topic you teach",
            "E.g., soil-moisture sensors → irrigation controller; "
            "pedometers → class dashboard",
            "Star a favorite from someone else",
        ],
        link="→ Open Shared Doc",
    )

    # 6. Reframe intro
    phase_intro_slide(
        prs, letter="R", title="Reframe",
        subtitle="Physical computing lives in the browser too — "
                 "the simulator is a real micro:bit",
        color=C_REFRAME, time_label="9:00",
    )

    # 7. 9:15 Breakout
    content_slide(
        prs, accent=C_REFRAME,
        do_label="YOU DO", phase_label="REFRAME",
        phase_color=C_REFRAME, time_label="9:15 (15 min)",
        title="Breakout: Teaching Hardware Without Hardware",
        body="In small groups:",
        bullets=[
            "#1 barrier to IoT in your building? "
            "(budget / shipping / breakage / absences / comfort)",
            "What does a simulator-first approach solve — and what do you lose?",
            "Post the trade-off in the shared doc",
        ],
    )

    # 8. Break #1
    break_slide(
        prs, time_label="9:30 (15 min)", title="Break #1",
        body="Stretch. Confirm MakeCode is loaded and at least one "
             "simulator is visible.",
    )

    # 9. Assemble intro
    phase_intro_slide(
        prs, letter="A", title="Assemble",
        subtitle="Sensor Node → Aggregator Node — "
                 "one conversation, two simulators",
        color=C_ASSEMBLE, time_label="9:45",
    )

    # 10. 9:45 Server Room Guardian (NEW)
    content_slide(
        prs, accent=C_ASSEMBLE,
        do_label="YOU DO", phase_label="ASSEMBLE · LIVE BUILD",
        phase_color=C_ASSEMBLE, time_label="9:45 (30 min)",
        title="Server Room Guardian",
        body="Watch a live build of sensor + aggregator, "
             "then swap in your own sensor.",
        bullets=[
            "Sensor broadcasts input.temperature() on radio group 7 every 2s",
            "Aggregator receives and displays — your hand warms the sim, "
            "the aggregator reacts",
            "Key move: LLM generates the code; "
            "toggle Blocks ↔ JavaScript to understand it",
            "Your turn: swap temp for light / accel / sound / compass — "
            "save TWO .hex files",
        ],
        link="→ makecode.microbit.org",
    )

    # 11. 10:15 Enhancements + Brainstorm (NEW)
    content_slide(
        prs, accent=C_ASSEMBLE,
        do_label="YOU DO", phase_label="ASSEMBLE · ENHANCE",
        phase_color=C_ASSEMBLE, time_label="10:15 (30 min)",
        title="Enhancements + Brainstorm",
        body="Two live upgrades, then you extend your own pair.",
        bullets=[
            "Liveness ping: aggregator sends PING, sensor replies ACK, "
            "dead-node warning on silence",
            "Data panel: wire incoming readings into MakeCode's live graph "
            "instead of scrolling LEDs",
            "Brainstorm: what's possible when nodes talk back + "
            "trends are visible?",
            "Build time: thresholds, alerts, bidirectional commands — "
            "log prompt → response → fix",
        ],
    )

    # 12. Break #2
    break_slide(
        prs, time_label="10:45 (15 min)", title="Break #2 — Show & Tell",
        body="Screenshot your paired simulators (sender + aggregator) "
             "in the chat.",
    )

    # 13. Level Up
    level_up_slide(prs)

    # 14. Fortify intro
    phase_intro_slide(
        prs, letter="F", title="Fortify",
        subtitle="Does your two-node system tell the truth?",
        color=C_FORTIFY, time_label="11:30",
    )

    # 15. 11:30 Experiment
    content_slide(
        prs, accent=C_FORTIFY,
        do_label="YOU DO", phase_label="FORTIFY",
        phase_color=C_FORTIFY, time_label="11:30 (15 min)",
        title="Experiment: Verify Your Simulated Data",
        body="Stress your Sensor + Aggregator:",
        bullets=[
            "Drag the sim's temp slider to extremes — "
            "does the running avg respond?",
            "Change the Sensor's radio group — Aggregator should fall silent",
            "Remove radio set group — watch the failure mode",
            "CtM: Task → Before → After → Takeaway",
            "Kit-arrives preview: real temp sensor reads CPU, "
            "not ambient (3–8°C hot)",
        ],
    )

    # 16. Transfer intro
    phase_intro_slide(
        prs, letter="T", title="Transfer",
        subtitle="Two-node lesson for YOUR classroom — then unlock your kit",
        color=C_TRANSFER, time_label="11:45",
    )

    # 17. 11:45 Build
    content_slide(
        prs, accent=C_TRANSFER,
        do_label="YOU DO", phase_label="TRANSFER",
        phase_color=C_TRANSFER, time_label="11:45 (15 min)",
        title="Build: Your Two-Node IoT Lesson",
        body="Customize the template:",
        bullets=[
            "What does the Sensor Node sense? Who holds it?",
            "What does the Aggregator Node decide / display? Who runs it?",
            "Student-facing CtM prompt · NGSS standard",
            "Post in shared doc + complete post-survey → "
            "kit ships this week",
        ],
        link="→ Open IoT Lesson Template",
    )

    # 18. Resource Kit
    resource_kit_slide(prs)

    # 19. Closing
    closing_slide(prs)

    out = "/home/user/wiki/static/pds/craft-pd-series/Workshop-3-Slides.pptx"
    prs.save(out)
    print(f"wrote {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
